"""Generate docs/BenchmarkingService.pptx — an overview + architecture deck.

Run with the project env plus python-pptx (no need to add it as a project dep):

    uv run --with python-pptx python docs/generate_pptx.py

Produces a 16-slide 16:9 deck: title, agenda, overview + key design decisions, the
design motif, two
architecture diagrams (the whole system, then inside the workload), the two-token auth
model, the benchmark catalog + run lifecycle, three benchmark slides, four results slides
(the 12-run matrix, what it measured, OpenShift vs KinD, AuthBridge overhead), and the
enact/report boundaries.

Content is sourced from docs/SERVICE_DESIGN_DECISIONS.md, docs/BENCHMARKS_PRIMER.md and the
generated results/12run-*.md documents. Every measured figure comes from our own runs -- the
benchmark slides over rows with intact telemetry, the results slides from the v1.24 matrices'
mirrored artifacts -- not from the benchmarks' published papers.

Slide titles carry their agenda number, so keep the agenda list and the title_band() prefixes
in step when adding or reordering slides.
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

# ---- palette ---------------------------------------------------------------
INK = RGBColor(0x1F, 0x2A, 0x37)       # near-black text
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NAVY = RGBColor(0x1B, 0x3A, 0x5C)      # title band
BLUE = RGBColor(0x2E, 0x6F, 0xB5)      # service
LTBLUE = RGBColor(0xE8, 0xF1, 0xFB)
CLIENT = RGBColor(0x5B, 0x8C, 0x5A)    # client (green)
LTGREEN = RGBColor(0xE9, 0xF3, 0xE8)
KC = RGBColor(0xB5, 0x5A, 0x2E)        # keycloak (rust)
LTORANGE = RGBColor(0xFB, 0xEE, 0xE4)
ROSSO = RGBColor(0x8E, 0x44, 0xAD)     # rossoctl (purple)
LTPURPLE = RGBColor(0xF1, 0xE8, 0xF6)
WORK = RGBColor(0x3D, 0x6E, 0x70)      # workload (teal)
LTTEAL = RGBColor(0xE4, 0xF0, 0xF0)
STORE = RGBColor(0x6B, 0x6B, 0x6B)     # mlflow/s3 (gray)
LTGRAY = RGBColor(0xEE, 0xEF, 0xF1)
BORDER = RGBColor(0xC7, 0xCE, 0xD6)
ACCENT = RGBColor(0xE8, 0x7A, 0x1E)    # arrow accent

prs = Presentation()
prs.slide_width = Emu(12192000)   # 13.333"
prs.slide_height = Emu(6858000)   # 7.5"
BLANK = prs.slide_layouts[6]

SW = prs.slide_width
SH = prs.slide_height


def _set_font(run, size, bold=False, color=INK, italic=False):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = "Calibri"
    run.font.color.rgb = color


def textbox(slide, x, y, w, h, lines, size=14, color=INK, align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP, bold=False):
    """lines: str, or list of (text, size, bold, color[, bullet_level])."""
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(lines, str):
        lines = [(lines, size, bold, color)]
    for i, spec in enumerate(lines):
        text, sz, bd, col = spec[0], spec[1], spec[2], spec[3]
        lvl = spec[4] if len(spec) > 4 else 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.level = lvl
        p.space_after = Pt(3)
        r = p.add_run()
        r.text = text
        _set_font(r, sz, bd, col)
    return tb


def box(slide, x, y, w, h, text, fill, line=BORDER, font=13, bold=True,
        font_color=INK, shape=MSO_SHAPE.ROUNDED_RECTANGLE, sub=None, sub_color=None):
    sp = slide.shapes.add_shape(shape, Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = line
    sp.line.width = Pt(1.0)
    sp.shadow.inherit = False
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    _set_font(r, font, bold, font_color)
    if sub:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = sub
        _set_font(r2, font - 3, False, sub_color or font_color)
    return sp


def connector(slide, x1, y1, x2, y2, color=ACCENT, width=2.0, dashed=False, arrow=True):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Emu(int(x1)), Emu(int(y1)),
                                    Emu(int(x2)), Emu(int(y2)))
    cn.line.color.rgb = color
    cn.line.width = Pt(width)
    ln = cn.line._get_or_add_ln()
    if arrow:
        tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
        ln.append(tail)
    if dashed:
        dash = ln.makeelement(qn("a:prstDash"), {"val": "dash"})
        ln.append(dash)
    cn.shadow.inherit = False
    return cn


def title_band(slide, title, subtitle=None):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Emu(950000))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    band.shadow.inherit = False
    tf = band.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(28)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    _set_font(r, 26, True, WHITE)
    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = subtitle
        _set_font(r2, 14, False, RGBColor(0xC9, 0xD9, 0xEC))
    return band


IN = 914400  # EMU per inch


def inch(v):
    return int(v * IN)



def grid(slide, x, y, w, h, rows, col_w, head_fill=NAVY, head_color=WHITE, font=10.5,
         first_col_bold=True):
    """A compact data table. python-pptx's native table beats hand-placed textboxes here: the
    column widths stay locked, so nothing drifts out of alignment when a label grows."""
    shape = slide.shapes.add_table(len(rows), len(rows[0]), Emu(int(x)), Emu(int(y)),
                                   Emu(int(w)), Emu(int(h)))
    tbl = shape.table
    tbl.first_row = True
    tbl.horz_banding = False
    for i, cw in enumerate(col_w):
        tbl.columns[i].width = Emu(int(cw))
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.margin_left = cell.margin_right = Pt(5)
            cell.margin_top = cell.margin_bottom = Pt(1)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = head_fill if r == 0 else (
                WHITE if r % 2 else RGBColor(0xF4, 0xF6, 0xF8))
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            run = para.add_run()
            run.text = str(val)
            _set_font(run, font, r == 0 or (c == 0 and first_col_bold),
                      head_color if r == 0 else INK)
    return tbl


def dot(cx, cy, n):
    """Numbered badge on the CURRENT slide -- reads the module-level `s`, so it must be called
    after `s` is rebound to the slide being drawn."""
    dia = 0.26 if len(str(n)) < 2 else 0.34  # widen for 2-digit numbers so they don't wrap
    d = s.shapes.add_shape(MSO_SHAPE.OVAL, int(cx) - inch(dia / 2), int(cy) - inch(dia / 2), inch(dia), inch(dia))
    d.fill.solid(); d.fill.fore_color.rgb = ACCENT; d.line.color.rgb = WHITE; d.line.width = Pt(1)
    d.shadow.inherit = False
    tf = d.text_frame; tf.word_wrap = False
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(n); _set_font(r, 10.5, True, WHITE)

# ============================================================ SLIDE 1: TITLE
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background(); bg.shadow.inherit = False
accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, inch(4.35), SW, inch(0.10))
accent.fill.solid(); accent.fill.fore_color.rgb = ACCENT; accent.line.fill.background(); accent.shadow.inherit = False
textbox(s, inch(0.9), inch(2.3), inch(11.5), inch(1.4),
        [("Benchmarking Service", 44, True, WHITE)], anchor=MSO_ANCHOR.MIDDLE)
textbox(s, inch(0.9), inch(3.5), inch(11.5), inch(0.8),
        [("Architecture & Design Overview", 22, False, RGBColor(0xC9, 0xD9, 0xEC))])
textbox(s, inch(0.9), inch(4.7), inch(11.5), inch(1.4),
        [("A pure-Python, HTTP-only service that deploys and evaluates agent benchmarks", 16, False, RGBColor(0xD8, 0xE3, 0xF0)),
         ("across multiple cluster-specific Rossoctl instances.", 16, False, RGBColor(0xD8, 0xE3, 0xF0))])

# ================================================================ SLIDE 2: AGENDA
s = prs.slides.add_slide(BLANK)
title_band(s, "Agenda", "What this deck covers, in order")
items = [
    ("Overview & key design decisions", "what the Service is, and the seven choices that shape it"),
    ("Auto-benchmarking design motif",
     "the shape of the system, before the wiring detail"),
    ("Architecture", "client, Service, and the per-cluster instance selected by JWT iss"),
    ("Architecture with workload specific components",
     "inside the workload: sidecar, user simulator, IBAC judge"),
    ("The two-token auth model", "why the caller's token is never forwarded upstream"),
    ("Benchmark catalog & run lifecycle", "the three benchmarks and the REST flow that drives them"),
    ("The three benchmarks — what they measure",
     "7.1 the difficulty ladder · 7.2 what each stresses · 7.3 the traps"),
    ("The canonical 12-run matrix",
     "8.1 what the 12 runs parameterize · 8.2 what they measured"),
    ("Cross-platform & plugin overhead",
     "9.1 OpenShift vs KinD, like-for-like · 9.2 what AuthBridge costs"),
    ("What the Service can & cannot enact", "the HTTP-only boundary, made explicit"),
]
# Two columns: a single column left ~54% of a 16:9 canvas empty. Split 4 + 3 and set the row pitch
# so the taller column reaches roughly the same depth as the content on the other slides.
SPLIT = 5  # 10 items, five per column
COLS = ((inch(0.75), inch(1.40), inch(4.95)), (inch(7.10), inch(7.75), inch(5.10)))
for idx, (head, sub) in enumerate(items):
    col = 0 if idx < SPLIT else 1
    row = idx if col == 0 else idx - SPLIT
    chip_x, text_x, text_w = COLS[col]
    y = inch(1.42) + row * inch(1.04)
    chip = s.shapes.add_shape(MSO_SHAPE.OVAL, chip_x, y + inch(0.06), inch(0.46), inch(0.46))
    chip.fill.solid(); chip.fill.fore_color.rgb = ACCENT
    chip.line.color.rgb = WHITE; chip.line.width = Pt(1.25); chip.shadow.inherit = False
    ctf = chip.text_frame; ctf.margin_left = ctf.margin_right = 0
    ctf.margin_top = ctf.margin_bottom = 0
    cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run(); cr.text = str(idx + 1); _set_font(cr, 14, True, WHITE)
    textbox(s, text_x, y, text_w, inch(0.32), [(head, 17, True, NAVY)])
    textbox(s, text_x, y + inch(0.33), text_w, inch(0.50),
            [(sub, 12.5, False, RGBColor(0x3A, 0x46, 0x54))])
# a hairline between the columns, so the split reads as deliberate structure
connector(s, inch(6.65), inch(1.40), inch(6.65), inch(6.05), color=BORDER, width=1.0, arrow=False)

# ================================================ SLIDE 3: OVERVIEW + DECISIONS
s = prs.slides.add_slide(BLANK)
title_band(s, "1.  Overview & Key Design Decisions")

# Left column: What it is
box(s, inch(0.45), inch(1.25), inch(5.9), inch(0.55), "What it is",
    NAVY, NAVY, font=16, font_color=WHITE, shape=MSO_SHAPE.RECTANGLE)
textbox(s, inch(0.5), inch(1.95), inch(5.85), inch(4.9),
        [("Deploys a benchmark's MCP tool + A2A agent, runs the evaluation, reads results, and exports them — all over HTTP.", 14, False, INK),
         ("Objectives", 14, True, BLUE),
         ("• Automate benchmarking lifecycle operations", 13.5, False, INK, 1),
         ("• Cross-user/cluster sharing of benchmark run results", 13.5, False, INK, 1),
         ("• Easy to configure in-cluster & cross-cluster benchmark runs", 13.5, False, INK, 1),
         ("• Asynchronous parallel benchmark runs", 13.5, False, INK, 1),
         ("• Secure, scalable, auditable & resilient", 13.5, False, INK, 1),
         ("3 benchmarks", 14, True, BLUE),
         ("• gsm8k — single-turn: one prompt in, one answer out, no dialogue partner", 13, False, INK, 1),
         ("• tau2 — multi-turn: a server-side user simulator LLM plays the customer", 13, False, INK, 1),
         ("• appworld — long-horizon: many API calls across simulated apps", 13, False, INK, 1),
         ("Client contract", 14, True, BLUE),
         ("• Point at a hostname, send Authorization: Bearer <caller JWT> — same from kind to prod", 13.5, False, INK, 1)])

# Right column: Key design decisions (the "overview chart")
box(s, inch(6.6), inch(1.25), inch(6.25), inch(0.55), "Key design decisions",
    NAVY, NAVY, font=16, font_color=WHITE, shape=MSO_SHAPE.RECTANGLE)
decisions = [
    ("Pure-Python, HTTPS-only to Rossoctl", "No shell-out: smaller image, no injection surface, structured errors, testable."),
    ("Two-token auth model", "Caller JWT attributes + routes only (never forwarded); Service mints its own benchmarker (ROPC) token to Rossoctl."),
    ("iss is the trust anchor", "The JWT issuer selects the per-instance config — no separate instance argument, so no confused-deputy escape."),
    ("Per-request ROPC login", "Fresh Service token every request — no expiry handling."),
    ("iss-keyed per-instance config", "One file per issuer: Rossoctl URL, benchmarker cred, MLflow/S3, optional Keycloak backchannel + workload route templates."),
    ("Enact only what the API allows", "Deploy/run/report/export — now including AuthBridge plugin presets (layer-3). Only workload Secrets remain out-of-band: prechecked and reported (424), never silently ignored."),
    ("MLflow with Service; S3 in the cloud", "MLflow is co-located with the Service (per-service traces it emits + reads), not in the workload cluster; S3 is an external cloud service — the shared cross-service sink."),
]
y = inch(1.95)
for head, body in decisions:
    b = box(s, inch(6.6), y, inch(6.25), inch(0.66),
            head, LTBLUE, BLUE, font=12.5, bold=True, font_color=INK,
            shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    b.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    b.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    b.text_frame.margin_left = Pt(8)
    p2 = b.text_frame.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run(); r2.text = body
    _set_font(r2, 10.5, False, RGBColor(0x3A, 0x46, 0x54))
    y += inch(0.70)


# ============================ SLIDE 4: AUTO-BENCHMARKING DESIGN MOTIF
# The conceptual shape, ahead of the wiring detail on the next slide: who talks to whom, and where
# results land. Deliberately omits Keycloak (auth has its own slide), the workload's internals (the
# slide after next), and anything cluster-specific.
s = prs.slides.add_slide(BLANK)
title_band(s, "2.  Auto-Benchmarking Design Motif",
           "One client, one Service, N workload deployments — and every result ends up in one place")

box(s, inch(0.4), inch(1.30), inch(2.8), inch(0.85),
    "Benchmarking Client", CLIENT, CLIENT, font=13.5, font_color=WHITE)

sv = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inch(0.4), inch(2.45), inch(2.8), inch(1.85))
sv.fill.solid(); sv.fill.fore_color.rgb = BLUE; sv.line.color.rgb = BLUE; sv.shadow.inherit = False
svtf = sv.text_frame; svtf.word_wrap = True; svtf.vertical_anchor = MSO_ANCHOR.TOP
svtf.margin_left = Pt(8); svtf.margin_top = Pt(7)
_p = svtf.paragraphs[0]; _p.alignment = PP_ALIGN.CENTER
_r = _p.add_run(); _r.text = "Benchmarking Service"; _set_font(_r, 13.5, True, WHITE)
for _txt in ("deploy the workload", "run the benchmark", "read + export results"):
    _pp = svtf.add_paragraph(); _pp.alignment = PP_ALIGN.LEFT; _pp.level = 1
    _rr = _pp.add_run(); _rr.text = "• " + _txt; _set_font(_rr, 11, False, WHITE)

box(s, inch(0.4), inch(4.75), inch(1.35), inch(0.95), "S3", STORE, ACCENT,
    font=13, font_color=WHITE, sub="shared sink", sub_color=LTGRAY)
box(s, inch(1.85), inch(4.75), inch(1.55), inch(0.95), "Experiment Tracker", STORE, STORE,
    font=11.5, font_color=WHITE)

# the instances container — stacked shadow implies "many of these"
CX3, CY3, CW3, CH3 = inch(4.05), inch(1.2), inch(8.75), inch(4.35)
for _off, _col in ((inch(0.16), RGBColor(0xF0, 0xF1, 0xF3)), (inch(0.08), RGBColor(0xE3, 0xE5, 0xE8))):
    _sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, CX3 + _off, CY3 + _off, CW3, CH3)
    _sh.fill.solid(); _sh.fill.fore_color.rgb = _col
    _sh.line.fill.background(); _sh.shadow.inherit = False
_c3 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, CX3, CY3, CW3, CH3)
_c3.fill.solid(); _c3.fill.fore_color.rgb = RGBColor(0xFB, 0xFC, 0xFD)
_c3.line.color.rgb = NAVY; _c3.line.width = Pt(1.5)
_c3.line._get_or_add_ln().append(
    _c3.line._get_or_add_ln().makeelement(qn("a:prstDash"), {"val": "dash"}))
_c3.shadow.inherit = False
textbox(s, inch(4.25), inch(1.30), inch(8.4), inch(0.4),
        [("Workload Deployment Instances", 13, True, NAVY)])

box(s, inch(4.35), inch(2.60), inch(2.55), inch(1.35),
    "Deployment helper for harnessed benchmark workload", ROSSO, ROSSO, font=12,
    font_color=WHITE)

_wg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inch(9.15), inch(2.30), inch(3.45), inch(1.60))
_wg.fill.solid(); _wg.fill.fore_color.rgb = LTTEAL
_wg.line.color.rgb = WORK; _wg.line.width = Pt(1.25); _wg.shadow.inherit = False
_wtf = _wg.text_frame; _wtf.vertical_anchor = MSO_ANCHOR.MIDDLE
_wp = _wtf.paragraphs[0]; _wp.alignment = PP_ALIGN.CENTER
_wr = _wp.add_run(); _wr.text = "Harnessed benchmark workload"; _set_font(_wr, 14, True, WORK)

box(s, inch(6.15), inch(4.72), inch(2.7), inch(0.70), "Telemetry Collector", WORK, WORK,
    font=12.5, font_color=WHITE)

connector(s, inch(1.8), inch(2.15), inch(1.8), inch(2.45)); dot(inch(2.05), inch(2.30), 1)
connector(s, inch(3.2), inch(3.15), inch(4.35), inch(3.15), color=ROSSO); dot(inch(3.78), inch(2.96), 2)
connector(s, inch(6.9), inch(3.00), inch(9.15), inch(3.00), color=ROSSO); dot(inch(8.02), inch(2.81), 3)
connector(s, inch(3.2), inch(4.10), inch(9.15), inch(3.70), color=WORK); dot(inch(7.70), inch(3.79), 4)
connector(s, inch(2.15), inch(4.30), inch(2.15), inch(4.75), color=STORE); dot(inch(2.15), inch(4.52), 5)
connector(s, inch(9.55), inch(3.90), inch(8.85), inch(4.72), color=WORK, dashed=True)
dot(inch(9.25), inch(4.33), "6a")
connector(s, inch(6.15), inch(5.07), inch(3.42), inch(5.22), color=STORE, dashed=True)
dot(inch(4.78), inch(5.145), "6b")
connector(s, inch(2.9), inch(4.75), inch(2.9), inch(4.30), color=STORE); dot(inch(2.9), inch(4.52), 7)
connector(s, inch(1.1), inch(4.30), inch(1.1), inch(4.75), color=ACCENT); dot(inch(0.88), inch(4.52), 8)

_lgl = ["1  Client → Service:  one bearer token, one hostname",
        "2  Service → deployment helper:  create the workload",
        "3  helper → harnessed workload:  MCP tool + A2A agent",
        "4  Service → harnessed workload:  run the benchmark"]
_lgr = ["5  Service → Experiment Tracker:  emit the per-task trace",
        "6a workload → Telemetry Collector  ·  6b collector → Tracker",
        "7  Experiment Tracker → Service:  read the records back",
        "8  Service → S3:  export the run's artifacts"]
_lb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inch(0.4), inch(5.85), inch(12.5), inch(1.15))
_lb.fill.solid(); _lb.fill.fore_color.rgb = RGBColor(0xF4, 0xF6, 0xF8)
_lb.line.color.rgb = BORDER; _lb.shadow.inherit = False
for _cx, _items in ((inch(0.6), _lgl), (inch(6.7), _lgr)):
    _ltf = s.shapes.add_textbox(_cx, inch(5.95), inch(6.0), inch(1.0)).text_frame
    _ltf.word_wrap = True
    for _i, _it in enumerate(_items):
        _pa = _ltf.paragraphs[0] if _i == 0 else _ltf.add_paragraph()
        _pa.space_after = Pt(0)
        _ru = _pa.add_run(); _ru.text = _it; _set_font(_ru, 11, False, INK)

# ================================================== SLIDE 5: ARCHITECTURE
s = prs.slides.add_slide(BLANK)
title_band(s, "3.  Architecture", "Relations between client, service, cluster-specific Keycloak / Rossoctl / workload, MLflow & S3")

# Client (off-cluster / host)
box(s, inch(0.4), inch(1.30), inch(2.8), inch(0.85),
    "Benchmarking Client", CLIENT, CLIENT, font=13.5, font_color=WHITE,
    sub="off-cluster (host / CI)", sub_color=RGBColor(0xE6, 0xF0, 0xE6))

# Service (center-left) — title anchored top, bullets below (no overlap)
sv = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inch(0.4), inch(2.45), inch(2.8), inch(1.85))
sv.fill.solid(); sv.fill.fore_color.rgb = BLUE; sv.line.color.rgb = BLUE; sv.shadow.inherit = False
svtf = sv.text_frame; svtf.word_wrap = True; svtf.vertical_anchor = MSO_ANCHOR.TOP
svtf.margin_left = Pt(8); svtf.margin_top = Pt(7)
p = svtf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Benchmarking Service"; _set_font(r, 13.5, True, WHITE)
for t in ("pure-Python · async", "iss-keyed InstanceRegistry", "per-request ROPC login",
          "deploy / run / report"):
    pp = svtf.add_paragraph(); pp.alignment = PP_ALIGN.LEFT; pp.level = 1
    rr = pp.add_run(); rr.text = "• " + t; _set_font(rr, 10, False, WHITE)

# MLflow + S3 sit on the SERVICE side, NOT inside the per-cluster workload instance:
# MLflow is co-located with the Service (same cluster); S3 is an external cloud service.
# S3 on the left, MLflow on the right (nearer the collector hop that feeds it).
box(s, inch(0.4), inch(4.75), inch(1.35), inch(0.95), "S3", STORE, ACCENT,
    font=13, font_color=WHITE, sub="cloud (external)", sub_color=LTGRAY)
ml = box(s, inch(1.85), inch(4.75), inch(1.35), inch(0.95), "MLflow", STORE, STORE,
         font=13, font_color=WHITE, sub="same cluster as Service", sub_color=LTGRAY)

# Per-cluster instance container (dashed) — stacked shadow implies "many instances"
CX, CY, CW, CH = inch(4.05), inch(1.2), inch(8.75), inch(4.35)
for off, col in ((inch(0.16), RGBColor(0xF0, 0xF1, 0xF3)), (inch(0.08), RGBColor(0xE3, 0xE5, 0xE8))):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, CX + off, CY + off, CW, CH)
    sh.fill.solid(); sh.fill.fore_color.rgb = col; sh.line.fill.background(); sh.shadow.inherit = False
cont = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, CX, CY, CW, CH)
cont.fill.solid(); cont.fill.fore_color.rgb = RGBColor(0xFB, 0xFC, 0xFD)
cont.line.color.rgb = NAVY; cont.line.width = Pt(1.5)
cont.line._get_or_add_ln().append(cont.line._get_or_add_ln().makeelement(qn("a:prstDash"), {"val": "dash"}))
cont.shadow.inherit = False
textbox(s, inch(4.2), inch(1.28), inch(8.4), inch(0.4),
        [("Per-cluster instance — selected by JWT  iss   (ykt3 · kind · …)", 12.5, True, NAVY)])

# Inside the cluster container
kc = box(s, inch(4.35), inch(1.95), inch(2.55), inch(1.0),
         "Keycloak", KC, KC, font=13, font_color=WHITE,
         sub="issuer + backchannel", sub_color=LTORANGE)
ro = box(s, inch(4.35), inch(3.2), inch(2.55), inch(1.1),
         "Rossoctl", ROSSO, ROSSO, font=13, font_color=WHITE,
         sub="backend API (server-side ops)", sub_color=LTPURPLE)

# Workload group
wg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inch(9.15), inch(1.95), inch(3.45), inch(2.35))
wg.fill.solid(); wg.fill.fore_color.rgb = LTTEAL; wg.line.color.rgb = WORK; wg.line.width = Pt(1.25)
wg.shadow.inherit = False
textbox(s, inch(9.25), inch(2.02), inch(3.25), inch(0.35),
        [("Benchmark Workload", 12, True, WORK)])
box(s, inch(9.35), inch(2.42), inch(3.05), inch(0.8), "MCP tool", WORK, WORK,
    font=12.5, font_color=WHITE, sub="exgentic-mcp-<benchmark>", sub_color=LTTEAL)
box(s, inch(9.35), inch(3.35), inch(3.05), inch(0.8), "A2A agent", WORK, WORK,
    font=12.5, font_color=WHITE, sub="exgentic-a2a-tool_calling-<b>", sub_color=LTTEAL)

# OTEL collector — lives in the workload cluster; forwards the agent's own OTLP spans to MLflow
# (the agent can't auth to MLflow directly). Optional / off by default (workload_otel).
box(s, inch(6.15), inch(4.72), inch(2.7), inch(0.70), "OTEL collector", WORK, WORK,
    font=12.5, font_color=WHITE, sub="forwards agent spans → MLflow", sub_color=LTTEAL)

# ---- connectors (numbered): dot() is a module-level helper (see above) ----

# 1 Client -> Keycloak (get caller JWT)
connector(s, inch(3.2), inch(1.7), inch(4.35), inch(2.25))
dot(inch(3.72), inch(1.95), 1)
# 2 Client -> Service (bearer)
connector(s, inch(1.8), inch(2.15), inch(1.8), inch(2.45))
dot(inch(2.05), inch(2.30), 2)
# 3 Service -> Keycloak (validate + mint token)
connector(s, inch(3.2), inch(2.75), inch(4.35), inch(2.6))
dot(inch(3.78), inch(2.66), 3)
# 4 Service -> Rossoctl
connector(s, inch(3.2), inch(3.55), inch(4.35), inch(3.65))
dot(inch(3.78), inch(3.58), 4)
# 5 Rossoctl -> Workload (deploys). 6 Service -> Workload (runs sessions).
# BOTH deliberately terminate on the GROUP boundary, not on an inner box: each addresses the MCP
# tool AND the A2A agent (5 creates both; 6 drives MCP sessions + A2A calls). Earlier they ended at
# y=2.9 / y=3.7, which lined up with the two inner boxes and so read as arrows aimed at a specific
# box and falling short. Aiming both at the group's vertical centre makes the group-level intent
# unambiguous; the per-component wiring is drawn on the next slide.
connector(s, inch(6.9), inch(3.55), inch(9.15), inch(3.22), color=ROSSO)
dot(inch(7.55), inch(3.455), 5)
# Arrow 6 must NOT take the direct line from the Service: (3.2,4.1) -> the group edge passes clean
# through the Rossoctl box and strikes out its label. Route it instead through the free 0.25" lane
# between Keycloak (bottom 2.95) and Rossoctl (top 3.20), then across to the group.
connector(s, inch(3.2), inch(4.05), inch(3.2), inch(3.07), color=WORK, arrow=False)
connector(s, inch(3.2), inch(3.07), inch(6.9), inch(3.07), color=WORK, arrow=False)
connector(s, inch(6.9), inch(3.07), inch(9.15), inch(3.35), color=WORK)
dot(inch(5.05), inch(3.07), 6)
# 7 Service -> MLflow: EMIT the Agent.Session trace (establishes it first, before the workload spans)
connector(s, inch(2.15), inch(4.30), inch(2.15), inch(4.75), color=STORE)
dot(inch(2.15), inch(4.52), 7)
# 8 is a TWO-HOP path (the agent cannot authenticate to MLflow itself), so both hops are badged
# 8a / 8b in flow order rather than leaving the first hop unlabelled, and both are dashed because
# the whole path is optional (workload_otel, off by default).
connector(s, inch(9.55), inch(4.15), inch(8.70), inch(4.72), color=WORK, dashed=True)
dot(inch(9.05), inch(4.48), "8a")
connector(s, inch(6.15), inch(5.07), inch(3.2), inch(5.24), color=STORE, dashed=True)
dot(inch(4.65), inch(5.155), "8b")
# 9 MLflow -> Service: READ back the Agent.Session traces (after the workload spans have landed)
connector(s, inch(2.9), inch(4.75), inch(2.9), inch(4.30), color=STORE)
dot(inch(2.9), inch(4.52), 9)
# 10 Service -> S3 (export) — terminal, after run completion
connector(s, inch(1.1), inch(4.30), inch(1.1), inch(4.75), color=ACCENT)
dot(inch(0.88), inch(4.52), 10)

# Legend — two-column panel below the frame
legend_l = [
    "1  Client logs in to Keycloak (ROPC) → caller JWT",
    "2  Client → Service:  Authorization: Bearer <caller JWT>",
    "3  Service validates JWT (JWKS) + mints its own",
    "      benchmarker token (ROPC, via backchannel)",
    "4  Service → Rossoctl: deploy / list agents + tools",
    "5  Rossoctl creates the MCP tool + A2A agent in-cluster",
]
legend_r = [
    "6  Service drives the MCP session + the A2A call itself;",
      "      the agent only issues execute_tool inside that session",
    "7  Service emits the Agent.Session trace → MLflow (first)",
    "8a A2A agent → OTEL collector   ·   8b collector → MLflow",
      "      (optional workload spans — off by default)",
    "9  Service reads back Agent.Session traces (MLflow)",
    "10  Service exports run.json / report.* (S3)",
]
lb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inch(0.4), inch(5.80), inch(12.5), inch(1.35))
lb.fill.solid(); lb.fill.fore_color.rgb = RGBColor(0xF4, 0xF6, 0xF8); lb.line.color.rgb = BORDER
lb.shadow.inherit = False
for col_x, items in ((inch(0.6), legend_l), (inch(6.7), legend_r)):
    ltf = s.shapes.add_textbox(col_x, inch(5.88), inch(6.0), inch(1.22)).text_frame
    ltf.word_wrap = True
    for i, item in enumerate(items):
        p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
        p.space_after = Pt(0)  # 7 lines in the right column; Pt(1) pressed line 10 into the border
        r = p.add_run(); r.text = item; _set_font(r, 10.5, False, INK)

# ============================ SLIDE 5: ARCHITECTURE — WORKLOAD SPECIFIC COMPONENTS
# Zooms into the "Benchmark Workload" group of the previous slide. The point of the slide is that
# the workload is NOT the same shape for every benchmark: two of the four components are optional,
# and how many LLMs a task involves depends on the benchmark and on the plugin preset.
s = prs.slides.add_slide(BLANK)
title_band(s, "4.  Architecture with Workload Specific Components",
           "Inside the Benchmark Workload — what every benchmark has, and what only some of them add")


def dash(sp):
    """Dashed outline = the component is conditional, not always deployed.

    The outline must CONTRAST with the fill or the dashes are invisible — the `box()` helper is
    normally called with line == fill for a flat look, which silently defeated this on the first
    render. White reads as a cut-out dash on every fill used here.
    """
    sp.line.color.rgb = WHITE
    sp.line.width = Pt(1.75)
    ln = sp.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    return sp


# Shared LLM gateway across the top: the agent, the tau2 user simulator and the IBAC judge all
# call it, so drawing it once as a bar keeps three flows short and non-crossing.
gw = box(s, inch(0.45), inch(1.20), inch(12.4), inch(0.6),
         "shared LLM gateway   ·   ete-litellm / litemaas   (external)", STORE, STORE,
         font=12.5, font_color=WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)

# Left column. The judge sits ABOVE the Service deliberately: the sidecar that calls it is in the
# lower row of the agent pod, so an upward-left arrow to the judge stays clear of the Service's own
# two arrows. Ordering these the other way forces the judge and Service flows to cross.
dash(box(s, inch(0.45), inch(2.35), inch(2.10), inch(0.75), "IBAC judge", KC, KC,
         font=11.5, font_color=WHITE, sub="an LLM call", sub_color=LTORANGE))
box(s, inch(0.45), inch(3.35), inch(2.10), inch(0.85), "Benchmarking Service", BLUE, BLUE,
    font=11, font_color=WHITE)

# The workload container. Its bottom band is intentionally deep (pods stop at 4.55) because the
# Service → MCP server flow is routed through it — the Service drives the session lifecycle on the
# MCP tool directly, it does not reach it via the agent.
CX2, CY2, CW2, CH2 = inch(2.72), inch(1.95), inch(10.13), inch(3.30)
cont2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, CX2, CY2, CW2, CH2)
cont2.fill.solid(); cont2.fill.fore_color.rgb = RGBColor(0xFB, 0xFC, 0xFD)
cont2.line.color.rgb = NAVY; cont2.line.width = Pt(1.5)
cont2.line._get_or_add_ln().append(
    cont2.line._get_or_add_ln().makeelement(qn("a:prstDash"), {"val": "dash"}))
cont2.shadow.inherit = False
# Label sits in the BAND BELOW the pods, not above them: the two vertical arrows to the LLM gateway
# leave the pods through the container's top edge, and a header there was struck through by arrow 2.
textbox(s, inch(2.90), inch(4.98), inch(9.7), inch(0.22),
        [("Benchmark Workload — one per benchmark, namespace team1", 11, True, NAVY)])

# --- A2A agent pod (left) ---
ap = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inch(3.00), inch(2.30), inch(4.55), inch(2.40))
ap.fill.solid(); ap.fill.fore_color.rgb = LTTEAL; ap.line.color.rgb = WORK; ap.line.width = Pt(1.25)
ap.shadow.inherit = False
textbox(s, inch(3.12), inch(2.32), inch(4.3), inch(0.24), [("A2A agent pod", 11, True, WORK)])
box(s, inch(3.12), inch(2.58), inch(4.31), inch(0.80), "agent container", WORK, WORK,
    font=12, font_color=WHITE, sub="exgentic-a2a-tool_calling-<b>  ·  the LLM loop", sub_color=LTTEAL)
dash(box(s, inch(3.12), inch(3.70), inch(4.31), inch(0.80), "AuthBridge sidecar", ROSSO, ROSSO,
         font=12, font_color=WHITE, sub="proxy  ·  only with --plugin-preset", sub_color=LTPURPLE))

# --- MCP tool pod (right) ---
mp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inch(8.20), inch(2.30), inch(4.50), inch(2.40))
mp.fill.solid(); mp.fill.fore_color.rgb = LTTEAL; mp.line.color.rgb = WORK; mp.line.width = Pt(1.25)
mp.shadow.inherit = False
textbox(s, inch(8.32), inch(2.32), inch(4.2), inch(0.24), [("MCP tool pod", 11, True, WORK)])
dash(box(s, inch(8.32), inch(2.58), inch(4.26), inch(0.80), "user simulator LLM", WORK, WORK,
         font=12, font_color=WHITE, sub="tau2 only  ·  plays the customer", sub_color=LTTEAL))
box(s, inch(8.32), inch(3.70), inch(4.26), inch(0.80), "MCP server", WORK, WORK,
    font=12, font_color=WHITE, sub="exgentic-mcp-<benchmark>  ·  tasks + evaluation", sub_color=LTTEAL)

# --- flows. Every arrow lands on the component it actually addresses, never on a group edge. ---
# 1 Service -> agent container (A2A send_prompt, once per task)
connector(s, inch(2.55), inch(3.42), inch(3.12), inch(3.12), color=BLUE)
dot(inch(2.84), inch(3.27), 1)
# 2 Service -> MCP server. Routed down the gap beside the container and along its bottom band: the
# Service calls create_session / evaluate_session / delete_session on the MCP tool ITSELF -- it does
# not reach the MCP tool via the agent.
connector(s, inch(2.64), inch(4.20), inch(2.64), inch(4.82), color=BLUE, arrow=False)
connector(s, inch(2.64), inch(4.82), inch(9.10), inch(4.82), color=BLUE, arrow=False)
connector(s, inch(9.10), inch(4.82), inch(9.10), inch(4.50), color=BLUE)
dot(inch(6.00), inch(4.82), 2)
# 3 agent -> LLM gateway (1 probe + N real chat calls)
connector(s, inch(5.25), inch(2.58), inch(5.25), inch(1.80), color=ACCENT)
dot(inch(5.25), inch(2.13), 3)
# 4 tau2 only: the user simulator in the MCP pod is a second LLM
connector(s, inch(10.45), inch(2.58), inch(10.45), inch(1.80), color=ACCENT)
dot(inch(10.45), inch(2.13), 4)
# 5 with a preset, the agent's MCP traffic goes through the sidecar first
# the badge is offset beside this arrow, not on it: the gap is only 0.32" and a centred badge
# covered the arrowhead, leaving the direction of travel unreadable.
connector(s, inch(5.25), inch(3.38), inch(5.25), inch(3.70), color=ROSSO)
dot(inch(5.62), inch(3.54), 5)
# 6 sidecar -> MCP server, once the action is authorized. The inter-pod corridor is 0.65" wide so
# this badge sits ON the line without touching either pod border.
connector(s, inch(7.43), inch(4.10), inch(8.32), inch(4.10), color=ROSSO)
dot(inch(7.88), inch(4.10), 6)
# 7 sidecar -> IBAC judge, per isAction tool call
connector(s, inch(3.12), inch(3.95), inch(2.55), inch(2.95), color=KC)
dot(inch(2.95), inch(3.63), 7)
# 8 the judge is itself an LLM call -- drawn, so the gateway's left third is not left unmoored
connector(s, inch(1.50), inch(2.35), inch(1.50), inch(1.80), color=KC, dashed=True)
dot(inch(1.50), inch(2.08), 8)

# --- bottom: per-benchmark matrix (left) + flow legend (right) ---
mtx = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inch(0.45), inch(5.40), inch(7.35), inch(1.88))
mtx.fill.solid(); mtx.fill.fore_color.rgb = RGBColor(0xF4, 0xF6, 0xF8)
mtx.line.color.rgb = BORDER; mtx.shadow.inherit = False
textbox(s, inch(0.62), inch(5.46), inch(7.0), inch(0.28),
        [("Which components a benchmark actually gets", 11.5, True, NAVY)])
COLS = ((inch(0.62), inch(1.15)), (inch(1.80), inch(2.50)),
        (inch(4.35), inch(2.15)), (inch(6.55), inch(1.15)))
ROWS = [
    # tool-call figures are MEDIANS per task, measured over the v1.24 matrices on both platforms
    # (gsm8k n=172 → 1, tau2 n=60 → 11, appworld n=41 → 13; appworld's spread is wide, 5–29).
    ("benchmark", "MCP tool image", "LLMs per task", "tool calls"),
    ("gsm8k", "exgentic-mcp-gsm8k", "1  (agent)", "~1"),
    ("tau2", "exgentic-mcp-tau2", "2  (+ user simulator)", "~11"),
    ("appworld", "exgentic-mcp-appworld", "1  (agent)", "~13"),
    ("+ any preset", "unchanged", "+1 IBAC judge call per action", "—"),
]
for ri, row in enumerate(ROWS):
    y = inch(5.80) + ri * inch(0.27)
    head = ri == 0
    for (cx, cw), cell in zip(COLS, row):
        textbox(s, cx, y, cw, inch(0.28),
                [(cell, 10, head, NAVY if head else INK)])

lg2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inch(8.05), inch(5.40), inch(4.8), inch(1.88))
lg2.fill.solid(); lg2.fill.fore_color.rgb = RGBColor(0xF4, 0xF6, 0xF8)
lg2.line.color.rgb = BORDER; lg2.shadow.inherit = False
flows = [
    "1  Service → agent:  send_prompt, once per task",
    "2  Service → MCP server:  list_tasks, create_session,",
    "      evaluate_session, delete_session  (the verdict)",
    "3  agent → gateway:  1 probe + N real chat calls",
    "4  tau2 only:  the user simulator is a 2nd LLM",
    "5  with a preset:  the agent's MCP calls go via the sidecar",
    "6  sidecar → MCP server, once authorized",
    "7  sidecar → judge, per isAction call",
    "8  the judge is itself an LLM call",
    "Dashed outline = present only in some configurations.",
    "No preset: 5 + 6 become one direct agent → MCP call.",
]
ftf = s.shapes.add_textbox(inch(8.22), inch(5.46), inch(4.5), inch(1.80)).text_frame
ftf.word_wrap = True
for i, item in enumerate(flows):
    p = ftf.paragraphs[0] if i == 0 else ftf.add_paragraph()
    p.space_after = Pt(0)  # 11 lines must fit the panel; Pt(1) overflowed past the slide edge
    r = p.add_run(); r.text = item; _set_font(r, 9.5, False, INK)

# ================================================ SLIDE 6: TWO-TOKEN AUTH
s = prs.slides.add_slide(BLANK)
title_band(s, "5.  The Two-Token Auth Model", "Why the caller's token is never forwarded upstream")

box(s, inch(0.5), inch(1.4), inch(5.9), inch(0.5), "Caller JWT  (inbound)", CLIENT, CLIENT,
    font=15, font_color=WHITE, shape=MSO_SHAPE.RECTANGLE)
textbox(s, inch(0.6), inch(2.0), inch(5.8), inch(3.6),
        [("Presented by the client as Authorization: Bearer.", 13.5, False, INK),
         ("• Used ONLY to attribute (preferred_username) and route (iss → instance)", 13, False, INK, 1),
         ("• Signature validated via the issuer's JWKS", 13, False, INK, 1),
         ("• aud / exp deliberately lenient in the dev/ops context", 13, False, INK, 1),
         ("• NEVER forwarded to Rossoctl", 13, True, KC, 1),
         ("benchmarker is special:", 13.5, True, CLIENT),
         ("• A caller JWT with preferred_username == benchmarker authorizes config ops (GET/PUT /config)", 13, False, INK, 1)])

box(s, inch(6.9), inch(1.4), inch(5.9), inch(0.5), "Service token  (outbound)", BLUE, BLUE,
    font=15, font_color=WHITE, shape=MSO_SHAPE.RECTANGLE)
textbox(s, inch(7.0), inch(2.0), inch(5.8), inch(3.6),
        [("Minted by the Service per request via ROPC (password grant) using the instance's benchmarker credential.", 13.5, False, INK),
         ("• Always fresh → no expiry handling", 13, False, INK, 1),
         ("• Presented to Rossoctl for all cluster-facing ops", 13, False, INK, 1),
         ("• The Service acts under its OWN identity, not the caller's", 13, True, BLUE, 1),
         ("Backchannel split (iss ≠ dialed host):", 13.5, True, BLUE),
         ("• When the iss host is unreachable (in-cluster kind), JWKS + token come from a configured backchannel URL; iss is matched, never dialed", 13, False, INK, 1)])

box(s, inch(1.6), inch(5.7), inch(10.1), inch(1.1),
    "Implication:  at the Rossoctl / cluster layer every action appears as the Service's identity — "
    "so per-user attribution & audit live in the Service, keyed on (iss, preferred_username).",
    LTBLUE, BLUE, font=13.5, bold=True, font_color=INK)

# ============================================ SLIDE 7: CATALOG + LIFECYCLE
s = prs.slides.add_slide(BLANK)
title_band(s, "6.  Benchmark Catalog & Run Lifecycle")

box(s, inch(0.45), inch(1.25), inch(5.75), inch(0.5), "Catalog (static registry)", NAVY, NAVY,
    font=15, font_color=WHITE, shape=MSO_SHAPE.RECTANGLE)
for i, (name, desc, col, lt) in enumerate([
    ("gsm8k", "single-turn · needs hf-secret + openai-secret", WORK, LTTEAL),
    ("tau2", "multi-turn · user-simulator LLM runs server-side in the MCP pod · raise timeout", ROSSO, LTPURPLE),
    ("appworld", "long-horizon across simulated apps · runs e2e · raise task_timeout_seconds", KC, LTORANGE),
]):
    y = inch(1.95) + i * inch(1.15)
    b = box(s, inch(0.45), y, inch(5.75), inch(1.0), name, lt, col, font=15, bold=True, font_color=col)
    b.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    b.text_frame.margin_left = Pt(10)
    p2 = b.text_frame.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run(); r2.text = desc; _set_font(r2, 11.5, False, INK)

box(s, inch(6.55), inch(1.25), inch(6.3), inch(0.5), "Lifecycle (REST)", NAVY, NAVY,
    font=15, font_color=WHITE, shape=MSO_SHAPE.RECTANGLE)
steps = [
    ("POST /benchmarks/{n}/deploy", "create MCP tool + A2A agent → 201"),
    ("GET /benchmarks/{n}/status", "poll until tool_ready & agent_ready"),
    ("POST /benchmarks/{n}/runs", "precheck → 202 run_id  (409 not deployed · 424 secret missing)"),
    ("GET …/runs/{id}", "poll pending → running → succeeded / failed"),
    ("GET …/runs/{id}/report", "MLflow records + artifacts  (409 if MLflow unset)"),
    ("download from S3", "run.json · report.* · token_report.* · span_report.* · manifest.json"),
]
y = inch(1.95)
for i, (ep, desc) in enumerate(steps):
    b = box(s, inch(6.9), y, inch(5.95), inch(0.66), ep, LTBLUE, BLUE, font=12.5,
            bold=True, font_color=INK)
    b.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    b.text_frame.margin_left = Pt(8)
    p2 = b.text_frame.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run(); r2.text = desc; _set_font(r2, 10.5, False, RGBColor(0x3A, 0x46, 0x54))
    if i < len(steps) - 1:
        connector(s, inch(6.75), y + inch(0.33), inch(6.9), y + inch(0.33) + inch(0.7),
                  color=BORDER, width=1.0, arrow=False)
    y += inch(0.75)
# down arrow spine
connector(s, inch(6.72), inch(2.1), inch(6.72), inch(6.4), color=ACCENT, width=1.5)


# ================================ SLIDES 8-10: THE THREE BENCHMARKS (from BENCHMARKS_PRIMER.md)
# Every figure is measured from our own v1.24/v1.23 runs, over rows with intact telemetry -- not
# quoted from the benchmarks' published papers.
s = prs.slides.add_slide(BLANK)
title_band(s, "7.1  The Three Benchmarks — a Difficulty Ladder",
           "Not interchangeable suites: each costs ~an order of magnitude more than the last")
grid(s, inch(0.45), inch(1.30), inch(12.4), inch(4.35), [
    ("", "gsm8k", "tau2", "appworld"),
    ("What it tests", "multi-step arithmetic", "multi-turn dialogue + tools", "long-horizon app automation"),
    ("Tasks measured", "172", "60 (50 clean)", "40 (35 clean)"),
    ("Pass rate", "0.98", "0.83", "0.00"),
    ("Input tokens / task", "343", "82,966", "210,792"),
    ("Output tokens / task", "205", "1,955", "21,499"),
    ("LLM calls / task", "2.1", "11.4", "23.6"),
    ("Tool calls / task", "1.1", "11.2", "13.4"),
    ("Median task latency", "5.3 s", "92 s", "232 s"),
    ("Slowest task seen", "62 s", "425 s", "581 s"),
    ("Model we use", "gpt-5-mini", "claude-sonnet-5", "gemini-2.5-pro"),
    ("Task pool", "8.5K (HuggingFace)", "114 (retail domain)", "grouped scenarios"),
], col_w=[inch(2.5), inch(3.3), inch(3.3), inch(3.3)], font=11)
_ban = box(s, inch(0.45), inch(6.05), inch(12.4), inch(0.95),
    "The scale gap is the headline: a tau2 task costs ~240x the input tokens of a gsm8k task, an "
    "appworld task ~615x. A 50-task gsm8k run is minutes; a 20-task appworld run is half an hour "
    "and millions of tokens. Budget by benchmark, not by task count.",
    LTGRAY, STORE, font=12.5, bold=True, font_color=INK)
_ban.text_frame.margin_left = _ban.text_frame.margin_right = Pt(18)

# ---- what each one is actually for ----
s = prs.slides.add_slide(BLANK)
title_band(s, "7.2  What Each Benchmark Stresses",
           "Why all three are in the matrix, and what a result from each does and does not tell you")
cards = [
    ("gsm8k", "the canary", WORK, LTTEAL, [
        "One prompt in, one answer out — no dialogue partner.",
        "~1 real LLM call + 1 tool call: the simplest agentic loop.",
        "Stresses almost nothing about the platform — which is the point.",
        "A failure here means INFRASTRUCTURE: deploy, auth, LLM reach, telemetry.",
        "Saturated at ~1.0, so it cannot discriminate models. Never read it as one.",
        "Deterministic enough that task 0 costs 320 input tokens on every cluster —"
        " we use that to prove two environments are comparable.",
    ]),
    ("tau2", "the discriminator", ROSSO, LTPURPLE, [
        "Multi-turn: a server-side USER SIMULATOR LLM plays the customer.",
        "Two models talk to each other, plus ~11 tool calls per task.",
        "Domain is retail (114 tasks) — the library default, never recorded in artifacts.",
        "The only leg where model choice dominates: 0.1 with gpt-5-mini vs 0.9 with"
        " claude-sonnet-5 on the SAME 10 tasks.",
        "Episodes are nondeterministic: across 6 samples of the same 10 tasks, 4 flip.",
        "At n=10 one task moves pass_rate by 0.10 — it cannot resolve less than that.",
    ]),
    ("appworld", "the stress test", KC, LTORANGE, [
        "Realistic chores across simulated apps: discover APIs, chain many calls.",
        "~24 LLM calls, ~13 tool calls, ~211k input, ~4 min per task.",
        "Evaluation is PROGRAMMATIC unit tests over final app state — no LLM judge.",
        "Pass rate 0.0 is honest, not broken: runs complete and tokens record —"
        " the agent just does not finish the job.",
        "~94% of tasks call finish: it believes it is done and the assertions"
        " disagree. No verification pass.",
        "Its job here is pipeline stress: long contexts, big traces, real timeouts.",
    ]),
]
x = inch(0.45)
for name, tag, col, lt, bullets in cards:
    hd = box(s, x, inch(1.30), inch(4.05), inch(0.62), f"{name}  —  {tag}", col, col,
             font=15, font_color=WHITE)
    body = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, inch(2.02), inch(4.05), inch(3.45))
    body.fill.solid(); body.fill.fore_color.rgb = lt
    body.line.color.rgb = col; body.line.width = Pt(1.0); body.shadow.inherit = False
    hd = box(s, x, inch(1.30), inch(3.95), inch(0.62), f"{name} \u2014 {tag}", col, col,
             font=15, font_color=WHITE)
    tf = body.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(10); tf.margin_top = Pt(30)
    # A shape's text frame defaults to MIDDLE anchoring and its FIRST paragraph inherits centred
    # alignment, so bullets floated in the middle of the box with line 1 centred and the rest left.
    # Both have to be set explicitly, and on every paragraph.
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, b in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        para.space_after = Pt(5)
        run = para.add_run(); run.text = "• " + b
        _set_font(run, 10, False, INK)
    x += inch(4.22)

grid(s, inch(0.45), inch(5.70), inch(12.4), inch(1.45), [
    ("If you want to \u2026", "use"),
    ("check a cluster / deploy / auth / telemetry path works", "gsm8k, 1\u201310 tasks"),
    ("exercise concurrency and volume cheaply", "gsm8k, 50 tasks at max_parallel_sessions=4"),
    ("compare models meaningfully", "tau2 \u2014 it discriminates; gsm8k saturates at ~1.0"),
    ("stress long contexts, long tasks, timeouts", "appworld"),
], col_w=[inch(7.4), inch(5.0)], font=11, first_col_bold=False)

# ---- the traps ----
s = prs.slides.add_slide(BLANK)
title_band(s, "7.3  Reading the Numbers — Five Things That Mislead",
           "Every one of these cost us a wrong conclusion first")
traps = [
    ("pass_rate = evaluated_pass / total",
     "A task that ERRORS before evaluation counts as not passed, so a low rate can mean "
     "\u201cfailed the task\u201d or \u201cnever got judged\u201d. Check the error column too."),
    ("The llm column overcounts real calls by one",
     "Every task issues an extra max_tokens=1 capability probe, counted as a chat span. "
     "llm=2 on gsm8k means ONE real call."),
    ("Implausibly small tokens = lost telemetry, and tokens == 0 will not catch it",
     "When a usage-bearing span is lost, what survives is the probe — whose own usage is 0/0 on "
     "reasoning models but 8/1 on claude-sonnet-5 and 1/0 on gemini. Use the structural test: "
     "llm <= 1 with tool >= 2 is impossible."),
    ("Output varies MORE than input, in most runs",
     "Measured OUT CV > IN CV in 27 of 33 legs. gsm8k's prompt is near-constant while answer "
     "length swings; only long-horizon appworld inverts it. Do not infer a direction — read the CV."),
    ("Task selection is deterministic",
     "A run takes the first max_tasks tasks, so the same task_id is the same task across runs "
     "and clusters, and a smaller run is a prefix of a larger one. That is what makes cross-platform "
     "comparison like-for-like — six legs matched to the byte."),
]
y = inch(1.32)
for i, (head, body_text) in enumerate(traps, 1):
    b = box(s, inch(0.45), y, inch(12.4), inch(1.00), f"{i}.  {head}", LTORANGE, KC,
            font=13, bold=True, font_color=INK)
    b.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    b.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    b.text_frame.margin_left = Pt(14); b.text_frame.margin_right = Pt(14)
    b.text_frame.margin_top = Pt(6)
    p2 = b.text_frame.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run(); r2.text = body_text
    _set_font(r2, 11, False, RGBColor(0x3A, 0x46, 0x54))
    y += inch(1.16)


# ============================ SLIDES 12-15: THE 12-RUN MATRIX AND WHAT IT SHOWED
# Every figure is read from the v1.24 matrices' own mirrored artifacts (results/12run-*.md).
s = prs.slides.add_slide(BLANK)
title_band(s, "8.1  The Canonical 12-Run Matrix",
           "One fixed set of 12 request bodies — the same on every platform, every version")
grid(s, inch(0.45), inch(1.30), inch(12.4), inch(3.95), [
    ("#", "benchmark", "tasks", "parallel", "what it varies"),
    ("1", "gsm8k", "1", "1", "baseline — the smoke test"),
    ("2", "gsm8k", "10", "1", "volume, still serial"),
    ("3", "gsm8k", "50", "4", "volume + concurrency"),
    ("4", "gsm8k", "5", "4", "model swap → Azure/gpt-4.1"),
    ("5", "gsm8k", "5", "4", "AuthBridge preset: auth-only"),
    ("6", "gsm8k", "5", "4", "AuthBridge preset: ibac-only"),
    ("7", "gsm8k", "5", "4", "AuthBridge preset: full (enforce)"),
    ("8", "gsm8k", "5", "4", "full + per-plugin override ibac:observe"),
    ("9", "tau2", "10", "1", "multi-turn + user simulator"),
    ("10", "tau2", "20", "4", "multi-turn under concurrency"),
    ("11", "appworld", "5", "1", "long-horizon, gemini-2.5-pro"),
    ("12", "appworld", "20", "4", "long-horizon under concurrency"),
], col_w=[inch(0.7), inch(1.7), inch(1.0), inch(1.2), inch(7.8)], font=10.5,
   first_col_bold=False)
_m = box(s, inch(0.45), inch(5.45), inch(12.4), inch(1.55),
    "Why a FIXED matrix: task selection is deterministic, so the same leg run anywhere executes the "
    "same tasks in the same order. That is what makes a difference attributable to the platform "
    "rather than to the workload — and it is why every leg deploys fresh, since reusing a warm agent "
    "silently loses telemetry. Driven by one command (reference/run-12.py, ~1h55m); the three "
    "generated documents in results/ derive every number from the mirrored artifacts.",
    LTGRAY, STORE, font=12, bold=False, font_color=INK)
_m.text_frame.margin_left = _m.text_frame.margin_right = Pt(18)

# ---- 7.2 what it measured ----
s = prs.slides.add_slide(BLANK)
title_band(s, "8.2  What the 12 Runs Measured",
           "v1.24 on OpenShift (Service on ykt3, workloads on ykt2) — 138 tasks, all 12 legs succeeded")
grid(s, inch(0.45), inch(1.30), inch(6.05), inch(4.15), [
    ("#", "bench", "pass", "wall", "input tokens"),
    ("1", "gsm8k", "1.00", "9 s", "320"),
    ("2", "gsm8k", "1.00", "42 s", "3,166"),
    ("3", "gsm8k", "1.00", "97 s", "15,684"),
    ("4", "gsm8k", "0.80", "26 s", "3,613"),
    ("5", "gsm8k", "1.00", "31 s", "1,564"),
    ("6", "gsm8k", "1.00", "35 s", "1,564"),
    ("7", "gsm8k", "1.00", "32 s", "1,564"),
    ("8", "gsm8k", "1.00", "30 s", "1,564"),
    ("9", "tau2", "0.70", "780 s", "852,240"),
    ("10", "tau2", "0.85", "425 s", "1,550,718"),
    ("11", "appworld", "0.00", "2164 s", "1,988,986"),
    ("12", "appworld", "0.00", "1769 s", "4,009,033"),
], col_w=[inch(0.6), inch(1.5), inch(1.0), inch(1.15), inch(1.8)], font=10.5,
   first_col_bold=False)
notes = [
    ("All eight gsm8k legs pass", "1.00 everywhere except #4, where ONE task of five fails on the "
     "gpt-4.1 swap. The four AuthBridge presets do not change the result — only the latency."),
    ("tau2 is the only leg that moves", "0.70 and 0.85. Across six samples of the same 10 tasks, four "
     "flip: at n=10 one task is worth 0.10, so this cannot resolve less than that."),
    ("appworld 0.00 is the honest result", "Runs complete, tokens record, evaluation says the goal "
     "was not met. ~94% of tasks self-declare finish."),
    ("Token attribution complete: 0 of 138 rows lost", "The previous matrix had 15 damaged rows that a "
     "tokens == 0 check reported as clean. Fresh deploy per leg, plus a structural detector."),
    ("8.4M input / 674k output tokens", "over 5,438 s of wall time — appworld alone is 71% of the "
     "input, on 25 of 138 tasks."),
]
y = inch(1.30)
for head, body_text in notes:
    b = box(s, inch(6.75), y, inch(6.10), inch(0.80), head, LTBLUE, BLUE, font=12,
            bold=True, font_color=INK)
    b.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    b.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    b.text_frame.margin_left = Pt(10); b.text_frame.margin_top = Pt(5)
    p2 = b.text_frame.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run(); r2.text = body_text
    _set_font(r2, 10, False, RGBColor(0x3A, 0x46, 0x54))
    y += inch(0.86)


# ---- 8.1 cross-platform ----
s = prs.slides.add_slide(BLANK)
title_band(s, "9.1  OpenShift vs KinD — Like-for-Like",
           "Same 12 request bodies, same Service version, verified-identical instance config")
grid(s, inch(0.45), inch(1.30), inch(7.55), inch(4.20), [
    ("#", "bench", "pass OCP", "pass KinD", "in OCP", "in KinD", ""),
    ("1", "gsm8k", "1.00", "1.00", "320", "320", "identical"),
    ("2", "gsm8k", "1.00", "1.00", "3,166", "3,166", "identical"),
    ("3", "gsm8k", "1.00", "1.00", "15,684", "15,684", "identical"),
    ("4", "gsm8k", "0.80", "0.80", "3,613", "4,184", ""),
    ("5", "gsm8k", "1.00", "1.00", "1,564", "1,564", "identical"),
    ("6", "gsm8k", "1.00", "1.00", "1,564", "1,564", "identical"),
    ("7", "gsm8k", "1.00", "1.00", "1,564", "1,564", "identical"),
    ("8", "gsm8k", "1.00", "1.00", "1,564", "1,564", "identical"),
    ("9", "tau2", "0.70", "0.70", "852,240", "847,530", ""),
    ("10", "tau2", "0.85", "0.80", "1,550,718", "1,552,641", ""),
    ("11", "appworld", "0.00", "0.00", "1,988,986", "1,150,541", ""),
    ("12", "appworld", "0.00", "0.00", "4,009,033", "3,646,292", ""),
], col_w=[inch(0.55), inch(1.25), inch(0.95), inch(1.0), inch(1.35), inch(1.35), inch(1.1)],
   font=10, first_col_bold=False)
find = [
    ("11 of 12 pass rates identical", "Only #10 differs, by one task of twenty (0.85 vs 0.80) — and "
     "tau2 is the leg we know flips run to run on either platform."),
    ("7 legs byte-identical on input tokens", "320 / 3,166 / 15,684 / 1,564 x4. Deterministic task "
     "selection plus the same model means identical work — the strongest available proof this is a "
     "like-for-like comparison, not merely a similar one."),
    ("Wall time is a dead heat", "5,438 s on OpenShift vs 5,429 s on a single KinD node — 0.2% apart "
     "in total, though individual legs differ by up to 8x."),
    ("0 lost-attribution rows on both sides", "138 OCP tasks, 135 KinD tasks."),
    ("So the platform is not the variable", "Where the two differ it traces to episode "
     "nondeterminism (tau2), small samples (#4 = one task), or appworld task timeouts — not to "
     "OpenShift vs KinD."),
]
y = inch(1.30)
for head, body_text in find:
    b = box(s, inch(8.25), y, inch(4.60), inch(0.80), head, LTTEAL, WORK, font=11.5,
            bold=True, font_color=INK)
    b.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    b.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    b.text_frame.margin_left = Pt(10); b.text_frame.margin_top = Pt(5)
    p2 = b.text_frame.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run(); r2.text = body_text
    _set_font(r2, 9.5, False, RGBColor(0x3A, 0x46, 0x54))
    y += inch(0.86)

# ---- 8.2 plugin overhead ----
s = prs.slides.add_slide(BLANK)
title_band(s, "9.2  What AuthBridge Costs",
           "Legs #3 and #5–8 ran byte-identical work, so latency differences are the plugin config")
grid(s, inch(0.45), inch(1.30), inch(6.35), inch(1.90), [
    ("component", "cost", "measured from"),
    ("Sidecar present at all", "~+9.5 s / task (OCP)", "flat across all four presets"),
    ("Tool call, sidecar but unjudged", "+0.50 s", "auth-only"),
    ("Tool call, judged", "+1.96 s on top", "full + observe (all 5 judged)"),
], col_w=[inch(2.6), inch(1.85), inch(1.90)], font=10.5)
grid(s, inch(0.45), inch(3.55), inch(6.35), inch(1.90), [
    ("", "OCP", "KinD", "ratio"),
    ("connect_mcp with sidecar", "3131 ms", "168 ms", "18.6x"),
    ("tool call, unjudged", "545 ms", "13 ms", "41.1x"),
    ("tool call, judged", "2507 ms", "1167 ms", "2.1x"),
], col_w=[inch(2.6), inch(1.25), inch(1.25), inch(1.25)], font=10.5)
pts = [
    ("The judge is itself an LLM call", "The sidecar POSTs to a chat-completions endpoint per "
     "authorized action, which is why it costs seconds and why its latency is so variable."),
    ("Most of the OCP figure is topology, not plugins",
     "The judged-call cost is comparable across platforms (2.1x) because it is dominated by a shared "
     "external gateway. The fixed cost is 8x apart because this OCP matrix is cross-cluster — "
     "token exchange over external routes. Do NOT quote +9.5 s as \u201cthe cost of AuthBridge\u201d; "
     "the single-node figures (~+1.2 s fixed, ~+1.1 s per judged call) are the better estimate."),
    ("Presets cannot be ranked from these numbers",
     "Only 1 of 6 IBAC legs authorized all of its tool calls: OCP judged 4/5, 3/5, 5/5 and KinD "
     "3/5, 4/5, 3/5. A leg that judged fewer calls shows a lower median — a mixture, not a saving. "
     "Whether that is decision caching under concurrency or a fail-open gap is still OPEN; the test "
     "is one ibac-only leg at parallelism 1."),
    ("Latency only", "The sidecar's CPU and memory cost is unmeasured — every infra field in these "
     "runs is 0.0."),
]
y = inch(1.30)
for head, body_text in pts:
    b = box(s, inch(7.05), y, inch(5.80), inch(1.30), head, LTPURPLE, ROSSO, font=11.5,
            bold=True, font_color=INK)
    b.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    b.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    b.text_frame.margin_left = Pt(10); b.text_frame.margin_top = Pt(5)
    p2 = b.text_frame.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run(); r2.text = body_text
    _set_font(r2, 9.5, False, RGBColor(0x3A, 0x46, 0x54))
    y += inch(1.36)

# ============================================ SLIDE 16: BOUNDARIES
s = prs.slides.add_slide(BLANK)
title_band(s, "10.  What the Service Can & Cannot Enact", "The HTTP-only boundary, made explicit")

box(s, inch(0.5), inch(1.35), inch(5.9), inch(0.5), "✓  Enacts over HTTP", CLIENT, CLIENT,
    font=15, font_color=WHITE, shape=MSO_SHAPE.RECTANGLE)
# Every line here is sized to fit on ONE line at its font size. A wrapped bullet returns to the
# column margin rather than hanging under its text, which reads as a broken line, so long items are
# split explicitly into a bullet + an indented continuation instead of being left to wrap.
textbox(s, inch(0.6), inch(2.0), inch(5.8), inch(4.6),
        [("• Deploy MCP tool + A2A agent (CPU/mem, image, env)", 13, False, INK, 1),
         ("• Deploy-time model swap (per-experiment agent)", 13, False, INK, 1),
         ("• authbridge_enabled → inject the sidecar (layer-2)", 13, False, INK, 1),
         ("• plugin_preset / plugins / on_error → AuthBridge layer-3", 13, False, INK, 1),
         ("→ pluginPreset/plugins/onError; operator renders ConfigMap", 11.5, False, RGBColor(0x3A, 0x46, 0x54), 2),
         ("• Run benchmark sessions; collect pass/fail + latency", 13, False, INK, 1),
         ("• Emit + read MLflow traces; export to S3", 13, False, INK, 1),
         ("• Service-owned config: MLflow + S3 via PUT /config", 13, False, INK, 1),
         ("benchmarker only", 11.5, False, RGBColor(0x3A, 0x46, 0x54), 2)])

box(s, inch(6.9), inch(1.35), inch(5.9), inch(0.5), "✗  Out-of-band (reports, doesn't do)", KC, KC,
    font=15, font_color=WHITE, shape=MSO_SHAPE.RECTANGLE)
textbox(s, inch(7.0), inch(2.0), inch(5.8), inch(4.6),
        [("• Cluster Secrets (hf-secret, openai-secret)", 13, False, INK, 1),
         ("operator provisions; run precheck returns 424 naming it", 11.5, False, RGBColor(0x3A, 0x46, 0x54), 2),
         ("• AuthBridge CLUSTER config — ibac judgeEndpoint / judgeModel", 13, False, INK, 1),
         ("in the platform-config ConfigMap, not the agent API", 11.5, False, RGBColor(0x3A, 0x46, 0x54), 2),
         ("• Any cluster-level API call — Rossoctl does these server-side", 13, False, INK, 1),
         ("Principle:", 13.5, True, KC),
         ("• Never silently ignore an un-enactable request —", 13, False, INK, 1),
         ("precheck (424) or reject (422), with an actionable reason", 11.5, False, RGBColor(0x3A, 0x46, 0x54), 2)])

box(s, inch(1.6), inch(5.75), inch(10.1), inch(0.75),
    "Cluster-agnostic by construction: works on kind / vanilla k8s / OpenShift; "
    "cross-cluster runs use per-instance route templates + a reachable internal issuer.",
    LTGRAY, STORE, font=12.5, bold=True, font_color=INK)

out = "docs/BenchmarkingService.pptx"
prs.save(out)
print("wrote", out, "with", len(prs.slides._sldIdLst), "slides")
